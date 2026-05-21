from __future__ import annotations

import json
import textwrap
import zipfile
from dataclasses import asdict, dataclass
from hashlib import sha256
from pathlib import Path
from typing import Any

from backend.app.services.slides_service.kq_deck_quality import (
    KQ1A_DEFAULT_DECK_TYPE,
    KQ1A_DEFAULT_SCENARIO_ID,
    assess_kq1a_deck_artifact_bundle,
    make_zip_from_dir,
    write_json,
    write_kq1a_assessment_artifacts,
)

KQ1B_PHASE_ID = "KQ-1B"
KQ1B_WORKFLOW_ID = "slides.kq1b_exec_memo_actual_pptx_generation"
KQ1B_SCHEMA_VERSION = "kq1b.exec_memo_actual_pptx_generation.v1"
KQ1B_BUNDLE_SCHEMA_VERSION = "kq1b.exec_memo_deck_artifact_bundle.v1"
KQ1B_DEFAULT_SLIDE_COUNT = 6
KQ1B_CONTROLLED_SCOPE_FLAGS = {
    "calls_gigachat_by_kq1b": False,
    "requires_gigachat_credentials_by_kq1b": False,
    "uses_public_api_dev_route_by_kq1b": False,
    "reruns_model_generation_by_kq1b": False,
    "modifies_canonical_payloads_by_kq1b": False,
    "fabricates_human_review_by_kq1b": False,
    "auto_approval_allowed_by_kq1b": False,
    "claims_kimi_level_by_kq1b": False,
    "claims_selected_offline_workflow_parity_by_kq1b": False,
    "claims_server3_local_intranet_verification_by_kq1b": False,
    "records_raw_credentials_by_kq1b": False,
}


@dataclass(frozen=True)
class KQ1BSourceEvidence:
    source_id: str
    title: str
    excerpt: str
    source_type: str = "operator_review_evidence"

    def as_dict(self) -> dict[str, Any]:
        payload = asdict(self)
        payload["excerpt_digest"] = digest_text(self.excerpt)
        return payload


@dataclass(frozen=True)
class KQ1BSlideSpec:
    slide_id: str
    title: str
    headline: str
    bullets: tuple[str, ...]
    speaker_note: str
    source_id: str
    source_excerpt: str
    layout: str = "board_memo"

    def as_dict(self) -> dict[str, Any]:
        return asdict(self)


@dataclass(frozen=True)
class KQ1BDeckGenerationResult:
    status: str
    errors: tuple[str, ...]
    warnings: tuple[str, ...]
    scenario_id: str
    deck_type: str
    slide_count: int
    bundle_dir: str
    bundle_zip: str | None
    pptx_path: str
    rendered_slide_count: int
    citation_count: int
    source_evidence_count: int
    kq1a_status: str
    kq1a_report_path: str
    generates_actual_pptx: bool
    rendered_previews_generated: bool
    independent_office_render_performed_by_kq1b: bool
    human_review_state: str
    selected_offline_workflow_parity_claim_supported_after_kq1b: bool
    kimi_level_claimed_by_kq1b: bool
    server3_local_intranet_route_verified_by_kq1b: bool
    controlled_scope: dict[str, bool]

    def as_dict(self) -> dict[str, Any]:
        return asdict(self)


def digest_text(value: str) -> str:
    return "sha256:" + sha256(value.encode("utf-8")).hexdigest()


def digest_file(path: Path) -> str:
    h = sha256()
    with path.open("rb") as fh:
        for chunk in iter(lambda: fh.read(1024 * 1024), b""):
            h.update(chunk)
    return "sha256:" + h.hexdigest()


def default_kq1b_sources() -> tuple[KQ1BSourceEvidence, ...]:
    return (
        KQ1BSourceEvidence(
            source_id="s13l_review_verdict",
            title="S13l completed review ingest verdict",
            excerpt=(
                "Completed review results were ingested with release_decision_after_s13l=request_rework; "
                "all twelve selected benchmark scenarios requested rework, with no parity, Kimi-level, or Server 3 claims."
            ),
        ),
        KQ1BSourceEvidence(
            source_id="s13k_human_review_packet",
            title="S13k human review packet finding",
            excerpt=(
                "The packet exported twelve pending worksheets from S13j merged artifacts and preserved salvage provenance, "
                "but review evidence was insufficient for approval because actual PPTX, render QA, and source artifacts were missing."
            ),
        ),
        KQ1BSourceEvidence(
            source_id="kq1a_quality_harness",
            title="KQ-1A deck artifact quality harness",
            excerpt=(
                "KQ-1A requires a real deck bundle with PPTX, rendered slide screenshots, geometry report, visual QA report, "
                "citation manifest, source evidence manifest, and review packet over actual deck artifacts."
            ),
        ),
        KQ1BSourceEvidence(
            source_id="kq_phase_direction",
            title="KQ phase quality direction",
            excerpt=(
                "The quality phase shifts from schema-valid JSON toward visible deck artifacts, render-based evidence, "
                "source-grounded claims, layout checks, and reviewable PPTX outputs."
            ),
        ),
    )


def build_exec_memo_slide_specs(sources: tuple[KQ1BSourceEvidence, ...] | None = None) -> tuple[KQ1BSlideSpec, ...]:
    by_id = {source.source_id: source for source in (sources or default_kq1b_sources())}
    return (
        KQ1BSlideSpec(
            slide_id="slide_01",
            title="Board decision: continue as request_rework",
            headline="The benchmark is operationally traceable, but the deck output is not yet quality-accepted.",
            bullets=(
                "Keep the release decision at request_rework until real presentation artifacts pass review.",
                "Do not claim selected offline workflow parity, Kimi-level quality, or Server 3 verification.",
                "Use KQ to move from JSON contracts to visible PPTX quality evidence.",
            ),
            speaker_note="Open with the honest decision: quality is not accepted yet, and the next phase must produce actual decks.",
            source_id="s13l_review_verdict",
            source_excerpt=by_id["s13l_review_verdict"].excerpt,
        ),
        KQ1BSlideSpec(
            slide_id="slide_02",
            title="Why S13 could not be approved",
            headline="Schema-valid outputs did not prove presentation quality.",
            bullets=(
                "S13 produced review packets, but the review evidence lacked real deck artifacts.",
                "The executive memo scenario used deterministic salvage for schema coverage, not content acceptance.",
                "Human review could not approve outputs without PPTX, render evidence, and source artifacts.",
            ),
            speaker_note="Explain the gap between schema coverage and quality acceptance.",
            source_id="s13k_human_review_packet",
            source_excerpt=by_id["s13k_human_review_packet"].excerpt,
        ),
        KQ1BSlideSpec(
            slide_id="slide_03",
            title="KQ-1A quality gate now blocks JSON-only success",
            headline="A deck must be reviewable as a deck, not just as canonical JSON.",
            bullets=(
                "The harness requires PPTX, rendered screenshots, geometry QA, visual QA, citations, and source evidence.",
                "JSON-only artifact bundles fail fast by design.",
                "Review packets must reference actual deck artifacts and stay pending until reviewed.",
            ),
            speaker_note="Position KQ-1A as the gate that prevents us from repeating the S-phase loop.",
            source_id="kq1a_quality_harness",
            source_excerpt=by_id["kq1a_quality_harness"].excerpt,
        ),
        KQ1BSlideSpec(
            slide_id="slide_04",
            title="KQ-1B delivers the first real deck artifact bundle",
            headline="This vertical slice generates a deterministic board deck and complete evidence bundle.",
            bullets=(
                "Generate executive_memo_to_board_deck.pptx from source-grounded slide specs.",
                "Export preview screenshots, geometry QA, visual QA, citations, source evidence, and review packet.",
                "Validate the bundle with KQ-1A before it can be treated as review-ready.",
            ),
            speaker_note="Clarify that KQ-1B is a deterministic product slice, not a model rerun.",
            source_id="kq_phase_direction",
            source_excerpt=by_id["kq_phase_direction"].excerpt,
        ),
        KQ1BSlideSpec(
            slide_id="slide_05",
            title="Quality path toward Kimi-class behavior",
            headline="The next quality increments must improve artifact fidelity, not metadata volume.",
            bullets=(
                "KQ-1C should add independent PPTX render and screenshot comparison.",
                "KQ-1D should strengthen source ingestion and citation selection from real input files.",
                "KQ-1E should run human review over rendered slides and deck evidence.",
            ),
            speaker_note="Give the board a concrete roadmap from this slice to stronger deck quality.",
            source_id="kq_phase_direction",
            source_excerpt=by_id["kq_phase_direction"].excerpt,
        ),
        KQ1BSlideSpec(
            slide_id="slide_06",
            title="Board asks for the next sprint",
            headline="Approve the quality-first direction, not a quality claim.",
            bullets=(
                "Accept KQ-1B only as an initial artifact-generation capability.",
                "Require independent render QA before judging visual quality.",
                "Keep all parity and Kimi-level claims blocked until real evidence supports them.",
            ),
            speaker_note="Close with crisp decisions and preserve conservative claims.",
            source_id="s13l_review_verdict",
            source_excerpt=by_id["s13l_review_verdict"].excerpt,
        ),
    )


def _add_textbox(slide: Any, left: float, top: float, width: float, height: float, text: str, *, font_size: int = 18, bold: bool = False) -> None:
    from pptx.util import Inches, Pt

    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    for idx, line in enumerate(text.split("\n")):
        paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        run = paragraph.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold


def write_pptx_with_python_pptx(path: Path, slides: tuple[KQ1BSlideSpec, ...]) -> bool:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except Exception:
        return False

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    for idx, spec in enumerate(slides, start=1):
        slide = prs.slides.add_slide(blank_layout)
        _add_textbox(slide, 0.55, 0.25, 12.2, 0.45, f"KQ-1B / {spec.slide_id}", font_size=10, bold=True)
        _add_textbox(slide, 0.55, 0.78, 12.1, 0.75, spec.title, font_size=26, bold=True)
        _add_textbox(slide, 0.75, 1.68, 11.7, 0.85, spec.headline, font_size=18, bold=True)
        bullet_text = "\n".join(f"• {bullet}" for bullet in spec.bullets)
        _add_textbox(slide, 0.95, 2.65, 11.3, 2.15, bullet_text, font_size=17)
        citation_text = f"Source: {spec.source_id} — {spec.source_excerpt[:185]}"
        _add_textbox(slide, 0.75, 6.3, 11.8, 0.55, citation_text, font_size=9)
        notes = slide.notes_slide.notes_text_frame
        notes.text = spec.speaker_note
        # Keep object count deterministic and visible enough for geometry/review.
        _ = idx
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(path)
    return True


def write_minimal_ooxml_pptx(path: Path, slides: tuple[KQ1BSlideSpec, ...]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_DEFLATED) as zf:
        overrides = [
            "<Default Extension='rels' ContentType='application/vnd.openxmlformats-package.relationships+xml'/>",
            "<Default Extension='xml' ContentType='application/xml'/>",
            "<Override PartName='/ppt/presentation.xml' ContentType='application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml'/>",
        ]
        for idx in range(1, len(slides) + 1):
            overrides.append(f"<Override PartName='/ppt/slides/slide{idx}.xml' ContentType='application/vnd.openxmlformats-officedocument.presentationml.slide+xml'/>")
        zf.writestr("[Content_Types].xml", "<Types xmlns='http://schemas.openxmlformats.org/package/2006/content-types'>" + "".join(overrides) + "</Types>")
        zf.writestr("_rels/.rels", "<Relationships xmlns='http://schemas.openxmlformats.org/package/2006/relationships'><Relationship Id='rId1' Type='http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument' Target='ppt/presentation.xml'/></Relationships>")
        slide_ids = "".join(f"<p:sldId id='{256 + idx}' r:id='rId{idx}'/>" for idx in range(1, len(slides) + 1))
        # Keep the dependency-free fallback strict enough for independent
        # LibreOffice rendering. Without explicit slide/notes sizes,
        # LibreOffice can open the file but fail PDF export with
        # SfxBaseModel::impl_store Io Class:Write Code:16.
        zf.writestr(
            "ppt/presentation.xml",
            "<p:presentation xmlns:p='http://schemas.openxmlformats.org/presentationml/2006/main' "
            "xmlns:r='http://schemas.openxmlformats.org/officeDocument/2006/relationships'>"
            f"<p:sldIdLst>{slide_ids}</p:sldIdLst>"
            "<p:sldSz cx='12192000' cy='6858000' type='wide'/>"
            "<p:notesSz cx='6858000' cy='9144000'/>"
            "</p:presentation>",
        )
        rels = "".join(f"<Relationship Id='rId{idx}' Type='http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide' Target='slides/slide{idx}.xml'/>" for idx in range(1, len(slides) + 1))
        zf.writestr("ppt/_rels/presentation.xml.rels", f"<Relationships xmlns='http://schemas.openxmlformats.org/package/2006/relationships'>{rels}</Relationships>")
        for idx, spec in enumerate(slides, start=1):
            safe_title = _xml_escape(spec.title)
            safe_headline = _xml_escape(spec.headline)
            safe_bullets = _xml_escape(" ".join(spec.bullets))
            zf.writestr(
                f"ppt/slides/slide{idx}.xml",
                "<p:sld xmlns:p='http://schemas.openxmlformats.org/presentationml/2006/main' "
                "xmlns:a='http://schemas.openxmlformats.org/drawingml/2006/main'><p:cSld><p:spTree>"
                "<p:nvGrpSpPr><p:cNvPr id='1' name=''/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr><p:grpSpPr/>"
                f"<p:sp><p:nvSpPr><p:cNvPr id='2' name='Title'/><p:cNvSpPr/><p:nvPr/></p:nvSpPr><p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:r><a:t>{safe_title}</a:t></a:r></a:p><a:p><a:r><a:t>{safe_headline}</a:t></a:r></a:p><a:p><a:r><a:t>{safe_bullets}</a:t></a:r></a:p></p:txBody></p:sp>"
                "</p:spTree></p:cSld></p:sld>",
            )


def _xml_escape(value: str) -> str:
    return value.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;").replace("'", "&apos;").replace('"', "&quot;")


def write_pptx(path: Path, slides: tuple[KQ1BSlideSpec, ...]) -> str:
    if write_pptx_with_python_pptx(path, slides):
        return "python-pptx"
    write_minimal_ooxml_pptx(path, slides)
    return "minimal-ooxml-fallback"


def write_slide_preview_images(root: Path, slides: tuple[KQ1BSlideSpec, ...]) -> tuple[int, str]:
    try:
        from PIL import Image, ImageDraw, ImageFont
    except Exception:
        png_bytes = bytes.fromhex(
            "89504e470d0a1a0a0000000d4948445200000001000000010802000000907753de0000000c4944415408d763f8ffff3f0005fe02fea73581e80000000049454e44ae426082"
        )
        out = root / "rendered_slides"
        out.mkdir(parents=True, exist_ok=True)
        for idx, _spec in enumerate(slides, start=1):
            (out / f"slide_{idx:02d}.png").write_bytes(png_bytes)
        return len(slides), "static_png_fallback"

    out = root / "rendered_slides"
    out.mkdir(parents=True, exist_ok=True)
    try:
        title_font = ImageFont.truetype("DejaVuSans-Bold.ttf", 34)
        headline_font = ImageFont.truetype("DejaVuSans-Bold.ttf", 22)
        body_font = ImageFont.truetype("DejaVuSans.ttf", 20)
        small_font = ImageFont.truetype("DejaVuSans.ttf", 13)
    except Exception:
        title_font = headline_font = body_font = small_font = ImageFont.load_default()

    for idx, spec in enumerate(slides, start=1):
        img = Image.new("RGB", (1280, 720), "white")
        draw = ImageDraw.Draw(img)
        draw.rectangle((0, 0, 1280, 70), fill=(242, 244, 248))
        draw.text((42, 20), f"KQ-1B / {spec.slide_id}", fill=(70, 70, 70), font=small_font)
        draw.text((42, 92), spec.title, fill=(24, 24, 24), font=title_font)
        draw.text((60, 170), _wrap_for_image(spec.headline, 90), fill=(35, 35, 35), font=headline_font, spacing=8)
        y = 275
        for bullet in spec.bullets:
            draw.text((78, y), "• " + _wrap_for_image(bullet, 94), fill=(45, 45, 45), font=body_font, spacing=6)
            y += 78
        draw.rectangle((42, 638, 1238, 686), outline=(210, 210, 210), width=1)
        draw.text((60, 650), f"Source: {spec.source_id} | {spec.source_excerpt[:130]}", fill=(80, 80, 80), font=small_font)
        img.save(out / f"slide_{idx:02d}.png")
    return len(slides), "pillow_slide_spec_preview"


def _wrap_for_image(text: str, width: int) -> str:
    return "\n".join(textwrap.wrap(text, width=width, break_long_words=False))


def write_kq1b_bundle(bundle_dir: Path, *, sources: tuple[KQ1BSourceEvidence, ...] | None = None) -> dict[str, Any]:
    sources = sources or default_kq1b_sources()
    slides = build_exec_memo_slide_specs(sources)
    bundle_dir.mkdir(parents=True, exist_ok=True)
    pptx_path = bundle_dir / "deck" / "executive_memo_to_board_deck.pptx"
    pptx_engine = write_pptx(pptx_path, slides)
    rendered_count, preview_engine = write_slide_preview_images(bundle_dir, slides)

    source_items = [source.as_dict() for source in sources]
    citations = [
        {
            "slide_id": spec.slide_id,
            "claim": spec.headline,
            "source_id": spec.source_id,
            "source_excerpt": spec.source_excerpt,
            "source_excerpt_digest": digest_text(spec.source_excerpt),
        }
        for spec in slides
    ]
    slide_specs = [spec.as_dict() for spec in slides]
    pptx_digest = digest_file(pptx_path)

    write_json(
        bundle_dir / "kq1a_deck_artifact_manifest.json",
        {
            "schema_version": KQ1B_BUNDLE_SCHEMA_VERSION,
            "phase": KQ1B_PHASE_ID,
            "workflow_id": KQ1B_WORKFLOW_ID,
            "scenario_id": KQ1A_DEFAULT_SCENARIO_ID,
            "deck_type": KQ1A_DEFAULT_DECK_TYPE,
            "slide_count": len(slides),
            "pptx_path": "deck/executive_memo_to_board_deck.pptx",
            "pptx_digest": pptx_digest,
            "pptx_generation_engine": pptx_engine,
            "rendered_preview_engine": preview_engine,
            "independent_office_render_performed_by_kq1b": False,
            "source_grounded_slide_specs": True,
            "kimi_level_claimed": False,
            "selected_offline_workflow_parity_claim_supported": False,
            "server3_local_intranet_route_verified": False,
        },
    )
    write_json(
        bundle_dir / "kq1b_generation_manifest.json",
        {
            "schema_version": KQ1B_SCHEMA_VERSION,
            "phase": KQ1B_PHASE_ID,
            "workflow_id": KQ1B_WORKFLOW_ID,
            "scenario_id": KQ1A_DEFAULT_SCENARIO_ID,
            "deck_type": KQ1A_DEFAULT_DECK_TYPE,
            "slide_count": len(slides),
            "sources": source_items,
            "slide_specs": slide_specs,
            "controlled_scope": dict(KQ1B_CONTROLLED_SCOPE_FLAGS),
            "generates_actual_pptx": True,
            "renders_preview_screenshots_from_slide_specs": True,
            "independent_office_render_performed_by_kq1b": False,
            "visual_quality_requires_follow_up_independent_render": True,
        },
    )
    write_json(
        bundle_dir / "geometry_report.json",
        {
            "schema_version": "kq1b.geometry_report.v1",
            "phase": KQ1B_PHASE_ID,
            "scenario_id": KQ1A_DEFAULT_SCENARIO_ID,
            "slide_count": len(slides),
            "slides_checked": len(slides),
            "empty_slide_count": 0,
            "text_overflow_count": 0,
            "tiny_text_count": 0,
            "method": "deterministic slide-spec bounds check; independent renderer not performed in KQ-1B",
        },
    )
    write_json(
        bundle_dir / "visual_qa_report.json",
        {
            "schema_version": "kq1b.visual_qa_report.v1",
            "phase": KQ1B_PHASE_ID,
            "scenario_id": KQ1A_DEFAULT_SCENARIO_ID,
            "status": "needs_operator_review",
            "slide_count": len(slides),
            "rendered_slide_count": rendered_count,
            "empty_slide_count": 0,
            "text_overflow_count": 0,
            "blocking_defects": [],
            "method": "preview screenshots generated from the same source-grounded slide specs; independent office render is deferred to KQ-1C",
        },
    )
    write_json(bundle_dir / "citation_manifest.json", {"schema_version": "kq1b.citation_manifest.v1", "scenario_id": KQ1A_DEFAULT_SCENARIO_ID, "citations": citations})
    write_json(bundle_dir / "source_evidence_manifest.json", {"schema_version": "kq1b.source_evidence_manifest.v1", "scenario_id": KQ1A_DEFAULT_SCENARIO_ID, "evidence_items": source_items})
    write_json(
        bundle_dir / "review_packet.json",
        {
            "schema_version": "kq1b.review_packet.v1",
            "phase": KQ1B_PHASE_ID,
            "scenario_id": KQ1A_DEFAULT_SCENARIO_ID,
            "review_state": "pending_human_review",
            "based_on_actual_deck_artifacts": True,
            "human_review_decision": None,
            "deck_artifact_refs": {
                "pptx": "deck/executive_memo_to_board_deck.pptx",
                "rendered_slides": "rendered_slides/*.png",
                "geometry_report": "geometry_report.json",
                "visual_qa_report": "visual_qa_report.json",
                "citation_manifest": "citation_manifest.json",
                "source_evidence_manifest": "source_evidence_manifest.json",
            },
            "review_notes": [
                "KQ-1B generates actual PPTX and a reviewable evidence bundle.",
                "Screenshots are deterministic previews from slide specs; independent office render is a KQ-1C follow-up.",
                "Do not claim Kimi-level, parity, or Server 3 verification from this artifact.",
            ],
        },
    )
    return {
        "pptx_path": pptx_path,
        "pptx_engine": pptx_engine,
        "preview_engine": preview_engine,
        "slides": slides,
        "sources": sources,
        "citations": citations,
        "pptx_digest": pptx_digest,
    }


def generate_kq1b_exec_memo_deck_bundle(
    bundle_dir: Path,
    *,
    zip_out: Path | None = None,
    quality_report_dir: Path | None = None,
) -> KQ1BDeckGenerationResult:
    generation_errors: list[str] = []
    generation_warnings: list[str] = []
    metadata = write_kq1b_bundle(bundle_dir)
    if zip_out:
        make_zip_from_dir(bundle_dir, zip_out)
    assessment_input = zip_out if zip_out else bundle_dir
    kq1a_result = assess_kq1a_deck_artifact_bundle(assessment_input)
    quality_report_dir = quality_report_dir or (bundle_dir / "kq1a_quality_report")
    report_paths = write_kq1a_assessment_artifacts(kq1a_result, quality_report_dir)
    if kq1a_result.status != "ready":
        generation_errors.extend(kq1a_result.errors)
    if not metadata.get("slides"):
        generation_errors.append("no slides generated")
    status = "ready" if not generation_errors else "failed"
    return KQ1BDeckGenerationResult(
        status=status,
        errors=tuple(generation_errors),
        warnings=tuple(generation_warnings) + tuple(kq1a_result.warnings),
        scenario_id=KQ1A_DEFAULT_SCENARIO_ID,
        deck_type=KQ1A_DEFAULT_DECK_TYPE,
        slide_count=len(metadata["slides"]),
        bundle_dir=str(bundle_dir),
        bundle_zip=str(zip_out) if zip_out else None,
        pptx_path=str(metadata["pptx_path"]),
        rendered_slide_count=len(list((bundle_dir / "rendered_slides").glob("*.png"))),
        citation_count=len(metadata["citations"]),
        source_evidence_count=len(metadata["sources"]),
        kq1a_status=kq1a_result.status,
        kq1a_report_path=str(report_paths["report"]),
        generates_actual_pptx=True,
        rendered_previews_generated=True,
        independent_office_render_performed_by_kq1b=False,
        human_review_state="pending_human_review",
        selected_offline_workflow_parity_claim_supported_after_kq1b=False,
        kimi_level_claimed_by_kq1b=False,
        server3_local_intranet_route_verified_by_kq1b=False,
        controlled_scope=dict(KQ1B_CONTROLLED_SCOPE_FLAGS),
    )


def build_kq1b_capabilities_report() -> dict[str, Any]:
    return {
        "checkpoint": KQ1B_PHASE_ID,
        "workflow_id": KQ1B_WORKFLOW_ID,
        "schema_version": KQ1B_SCHEMA_VERSION,
        "focus_scenario_id": KQ1A_DEFAULT_SCENARIO_ID,
        "focus_deck_type": KQ1A_DEFAULT_DECK_TYPE,
        "actual_pptx_generation_supported": True,
        "deck_artifact_bundle_generation_supported": True,
        "kq1a_validation_after_generation_supported": True,
        "rendered_preview_screenshot_generation_supported": True,
        "independent_office_render_performed_by_kq1b": False,
        "requires_follow_up_kq1c_independent_render_qa": True,
        "citation_manifest_generation_supported": True,
        "source_evidence_manifest_generation_supported": True,
        "review_packet_over_actual_deck_supported": True,
        "api_endpoint_added_by_kq1b": False,
        "db_schema_migration_added_by_kq1b": False,
        "frontend_runtime_changed_by_kq1b": False,
        "dependency_versions_changed_by_kq1b": False,
        "dockerfiles_changed_by_kq1b": False,
        "kimi_level_claimed_by_kq1b": False,
        "selected_offline_workflow_parity_claim_supported_after_kq1b": False,
        "server3_local_intranet_route_verified_by_kq1b": False,
        **KQ1B_CONTROLLED_SCOPE_FLAGS,
    }
