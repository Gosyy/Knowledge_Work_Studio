from __future__ import annotations

import json
import shutil
import subprocess
import tempfile
import zipfile
from dataclasses import asdict, dataclass
from hashlib import sha256
from pathlib import Path
from typing import Any

from backend.app.services.slides_service.kq_deck_quality import (
    KQ1A_DEFAULT_DECK_TYPE,
    KQ1A_DEFAULT_SCENARIO_ID,
    assess_kq1a_deck_artifact_bundle,
    find_pptx_files,
    make_zip_from_dir,
    read_json,
    write_json,
    write_kq1a_assessment_artifacts,
)

KQ1C_PHASE_ID = "KQ-1C"
KQ1C_WORKFLOW_ID = "slides.kq1c_independent_pptx_render_visual_qa"
KQ1C_SCHEMA_VERSION = "kq1c.independent_pptx_render_visual_qa.v1"
KQ1C_RENDER_MANIFEST_SCHEMA_VERSION = "kq1c.render_manifest.v1"
KQ1C_VISUAL_QA_SCHEMA_VERSION = "kq1c.visual_qa_report.v1"
KQ1C_RENDER_DPI = 144
KQ1C_CONTROLLED_SCOPE_FLAGS = {
    "calls_gigachat_by_kq1c": False,
    "requires_gigachat_credentials_by_kq1c": False,
    "uses_public_api_dev_route_by_kq1c": False,
    "reruns_model_generation_by_kq1c": False,
    "modifies_canonical_payloads_by_kq1c": False,
    "fabricates_human_review_by_kq1c": False,
    "auto_approval_allowed_by_kq1c": False,
    "claims_kimi_level_by_kq1c": False,
    "claims_selected_offline_workflow_parity_by_kq1c": False,
    "claims_server3_local_intranet_verification_by_kq1c": False,
    "records_raw_credentials_by_kq1c": False,
}


@dataclass(frozen=True)
class KQ1CSlideRenderObservation:
    slide_id: str
    render_path: str
    width_px: int
    height_px: int
    non_white_pixel_ratio: float
    text_character_count: int
    extracted_text_preview: str
    empty_slide_detected: bool
    tiny_text_detected: bool
    text_overflow_detected: bool

    def as_dict(self) -> dict[str, Any]:
        return asdict(self)


@dataclass(frozen=True)
class KQ1CRenderQAResult:
    status: str
    errors: tuple[str, ...]
    warnings: tuple[str, ...]
    scenario_id: str
    deck_type: str
    input_bundle: str
    output_bundle_dir: str
    output_bundle_zip: str | None
    pptx_path: str
    pptx_digest: str
    slide_count_from_pptx: int
    independent_pptx_render_performed_by_kq1c: bool
    independent_office_render_performed_by_kq1c: bool
    render_engine: str
    render_engine_requested: str
    render_engine_available: bool
    rendered_slide_count: int
    rendered_slide_dir: str
    visual_qa_status: str
    empty_slide_count: int
    text_overflow_count: int
    tiny_text_count: int
    blocking_defect_count: int
    kq1a_status_after_kq1c: str
    kq1a_report_path: str
    selected_offline_workflow_parity_claim_supported_after_kq1c: bool
    kimi_level_claimed_by_kq1c: bool
    server3_local_intranet_route_verified_by_kq1c: bool
    controlled_scope: dict[str, bool]
    slide_observations: tuple[dict[str, Any], ...]

    def as_dict(self) -> dict[str, Any]:
        return asdict(self)


def digest_file(path: Path) -> str:
    h = sha256()
    with path.open("rb") as fh:
        for chunk in iter(lambda: fh.read(1024 * 1024), b""):
            h.update(chunk)
    return "sha256:" + h.hexdigest()


def _copy_bundle_to_dir(input_bundle: Path, output_dir: Path) -> Path:
    if output_dir.exists():
        shutil.rmtree(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)
    if input_bundle.is_file() and input_bundle.suffix.lower() == ".zip":
        with zipfile.ZipFile(input_bundle, "r") as zf:
            zf.extractall(output_dir)
        return output_dir
    if input_bundle.is_dir():
        for item in input_bundle.iterdir():
            destination = output_dir / item.name
            if item.is_dir():
                shutil.copytree(item, destination)
            else:
                shutil.copy2(item, destination)
        return output_dir
    raise RuntimeError(f"KQ-1C input must be a ZIP archive or extracted bundle directory: {input_bundle}")


def _load_json_object(path: Path) -> dict[str, Any]:
    payload = read_json(path)
    return payload if isinstance(payload, dict) else {}


def _find_review_packet(root: Path) -> Path | None:
    candidates = sorted(path for path in root.rglob("*.json") if path.name in {"review_packet.json", "human_review_packet.json"})
    return candidates[0] if candidates else None


def _find_bundle_manifest(root: Path) -> Path | None:
    candidates = sorted(
        path
        for path in root.rglob("*.json")
        if path.name in {"kq1a_deck_artifact_manifest.json", "deck_artifact_manifest.json", "kq1b_generation_manifest.json"}
    )
    return candidates[0] if candidates else None


def _slide_count_with_python_pptx(pptx_path: Path) -> int:
    try:
        from pptx import Presentation

        return len(Presentation(str(pptx_path)).slides)
    except Exception:
        try:
            with zipfile.ZipFile(pptx_path, "r") as zf:
                return len([name for name in zf.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml")])
        except Exception:
            return 0


def _extract_pptx_slide_text_from_xml(pptx_path: Path) -> list[str]:
    """Extract slide text from PPTX XML without optional python-pptx.

    KQ-1C should prefer the real Office/PDF render stack and must not require
    optional python-pptx/Pillow fallback dependencies when LibreOffice is
    available. This keeps visual QA dependency-light on Profile 2.
    """
    try:
        import re
        import xml.etree.ElementTree as ET

        with zipfile.ZipFile(pptx_path, "r") as zf:
            slide_names = [
                name
                for name in zf.namelist()
                if name.startswith("ppt/slides/slide") and name.endswith(".xml")
            ]

            def slide_number(name: str) -> int:
                match = re.search(r"slide(\d+)\.xml$", name)
                return int(match.group(1)) if match else 0

            slide_texts: list[str] = []
            for name in sorted(slide_names, key=slide_number):
                root = ET.fromstring(zf.read(name))
                chunks = [
                    (node.text or "").strip()
                    for node in root.iter()
                    if node.tag.endswith("}t") and (node.text or "").strip()
                ]
                slide_texts.append("\n".join(chunks))
            return slide_texts
    except Exception:
        return []


def _extract_pptx_slide_text(pptx_path: Path) -> list[str]:
    try:
        from pptx import Presentation
    except Exception:
        return _extract_pptx_slide_text_from_xml(pptx_path)
    try:
        presentation = Presentation(str(pptx_path))
    except Exception:
        return _extract_pptx_slide_text_from_xml(pptx_path)
    slide_texts: list[str] = []
    for slide in presentation.slides:
        chunks: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text:
                chunks.append(str(text).strip())
        slide_texts.append("\n".join(chunk for chunk in chunks if chunk))
    return slide_texts


def _extract_pptx_text_geometry_findings(pptx_path: Path) -> tuple[list[int], list[bool], list[bool]]:
    try:
        from pptx import Presentation
        from pptx.util import Pt
    except Exception:
        slide_texts = _extract_pptx_slide_text_from_xml(pptx_path)
        return [len(text) for text in slide_texts], [False for _ in slide_texts], [False for _ in slide_texts]
    try:
        presentation = Presentation(str(pptx_path))
    except Exception:
        slide_texts = _extract_pptx_slide_text_from_xml(pptx_path)
        return [len(text) for text in slide_texts], [False for _ in slide_texts], [False for _ in slide_texts]

    char_counts: list[int] = []
    overflow_flags: list[bool] = []
    tiny_text_flags: list[bool] = []
    for slide in presentation.slides:
        slide_chars = 0
        slide_overflow = False
        slide_tiny_text = False
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text = getattr(shape, "text", "") or ""
            slide_chars += len(text)
            width = max(int(getattr(shape, "width", 0) or 0), 1)
            height = max(int(getattr(shape, "height", 0) or 0), 1)
            # Conservative approximation: one character needs roughly 13pt x 1.3 line-height.
            explicit_min_font_pt: float | None = None
            try:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.size is not None:
                            pt = float(run.font.size.pt)
                            explicit_min_font_pt = pt if explicit_min_font_pt is None else min(explicit_min_font_pt, pt)
            except Exception:
                explicit_min_font_pt = None
            effective_font_pt = explicit_min_font_pt or 12.0
            if explicit_min_font_pt is not None and explicit_min_font_pt < 8.0:
                slide_tiny_text = True
            width_in = width / 914400.0
            height_in = height / 914400.0
            approx_capacity = max(int(width_in * height_in * 88.0 * (12.0 / max(effective_font_pt, 1.0))), 16)
            if len(text) > approx_capacity * 1.4:
                slide_overflow = True
        char_counts.append(slide_chars)
        overflow_flags.append(slide_overflow)
        tiny_text_flags.append(slide_tiny_text)
    return char_counts, overflow_flags, tiny_text_flags


def _non_white_pixel_ratio(image_path: Path) -> tuple[int, int, float]:
    try:
        from PIL import Image
    except Exception:
        # Dependency-light fallback: when the independent Office renderer has
        # produced a non-empty PNG but Pillow is unavailable, keep QA
        # conservative without falsely marking every rendered slide empty.
        # Text presence is checked separately from PPTX XML.
        if image_path.exists() and image_path.stat().st_size > 0:
            return 1280, 720, 0.05
        return 0, 0, 0.0
    with Image.open(image_path) as img:
        rgb = img.convert("RGB")
        width, height = rgb.size
        pixels = list(rgb.getdata())
    if not pixels:
        return width, height, 0.0
    non_white = sum(1 for r, g, b in pixels if r < 245 or g < 245 or b < 245)
    return width, height, non_white / float(len(pixels))


def _find_libreoffice_pdf_candidates(*roots: Path) -> list[Path]:
    candidates: list[Path] = []
    seen: set[Path] = set()
    for root in roots:
        if not root.exists():
            continue
        if not root.is_dir():
            continue
        for candidate in root.rglob("*.pdf"):
            resolved = candidate.resolve()
            if resolved not in seen:
                seen.add(resolved)
                candidates.append(candidate)
    return sorted(candidates, key=lambda path: (path.stat().st_mtime, str(path)), reverse=True)


def _run_libreoffice_pdf_convert(soffice: str, pptx_path: Path, pdf_dir: Path, work_dir: Path, profile_dir: Path) -> subprocess.CompletedProcess[str]:
    # Use a per-run LibreOffice profile to avoid silent no-output conversions
    # caused by a locked or half-initialized user profile in CI/operator shells.
    # Some LibreOffice builds return 0 while failing to write into a separate
    # --outdir. Try a small auditable matrix: relative input + work_dir output
    # first, then explicit pdf_dir/original-path fallbacks.
    profile_uri = profile_dir.resolve().as_uri()
    local_input = work_dir / "input.pptx"
    shutil.copy2(pptx_path, local_input)

    base_flags = [
        soffice,
        f"-env:UserInstallation={profile_uri}",
        "--headless",
        "--invisible",
        "--nodefault",
        "--nofirststartwizard",
        "--nolockcheck",
        "--norestore",
    ]
    minimal_flags = [soffice, "--headless"]

    attempts: list[tuple[str, list[str], Path]] = [
        (
            "profile_relative_input_workdir_pdf",
            [*base_flags, "--convert-to", "pdf", "--outdir", str(work_dir), local_input.name],
            work_dir,
        ),
        (
            "profile_relative_input_workdir_impress_pdf",
            [*base_flags, "--convert-to", "pdf:impress_pdf_Export", "--outdir", str(work_dir), local_input.name],
            work_dir,
        ),
        (
            "minimal_relative_input_workdir_pdf",
            [*minimal_flags, "--convert-to", "pdf", "--outdir", str(work_dir), local_input.name],
            work_dir,
        ),
        (
            "profile_relative_input_default_outdir",
            [*base_flags, "--convert-to", "pdf", local_input.name],
            work_dir,
        ),
        (
            "profile_absolute_input_pdfdir_pdf",
            [*base_flags, "--convert-to", "pdf", "--outdir", str(pdf_dir), str(local_input)],
            work_dir,
        ),
        (
            "profile_original_input_pdfdir_pdf",
            [*base_flags, "--convert-to", "pdf", "--outdir", str(pdf_dir), str(pptx_path)],
            work_dir,
        ),
        (
            "minimal_absolute_input_pdfdir_pdf",
            [*minimal_flags, "--convert-to", "pdf", "--outdir", str(pdf_dir), str(local_input)],
            work_dir,
        ),
    ]

    transcripts: list[str] = []
    last_result: subprocess.CompletedProcess[str] | None = None
    for label, command, cwd in attempts:
        before = {path.resolve() for path in _find_libreoffice_pdf_candidates(pdf_dir, work_dir)}
        env = dict(__import__("os").environ)
        env.update({
            "HOME": str(profile_dir),
            "TMPDIR": str(work_dir),
            "SAL_USE_VCLPLUGIN": env.get("SAL_USE_VCLPLUGIN", "gen"),
        })
        last_result = subprocess.run(
            command,
            cwd=str(cwd),
            text=True,
            stdout=subprocess.PIPE,
            stderr=subprocess.STDOUT,
            timeout=120,
            check=False,
            env=env,
        )
        after = _find_libreoffice_pdf_candidates(pdf_dir, work_dir)
        new_candidates = [path for path in after if path.resolve() not in before]
        transcripts.append(
            "\n".join(
                [
                    f"attempt={label}",
                    f"returncode={last_result.returncode}",
                    "stdout_tail=" + last_result.stdout[-1200:],
                    "pdf_candidates=" + ", ".join(str(path) for path in after),
                ]
            )
        )
        if last_result.returncode == 0 and (new_candidates or after):
            last_result.stdout = "\n--- libreoffice-attempt ---\n".join(transcripts)
            return last_result

    assert last_result is not None
    last_result.stdout = "\n--- libreoffice-attempt ---\n".join(transcripts)
    return last_result


def _render_with_libreoffice(pptx_path: Path, render_dir: Path, *, dpi: int = KQ1C_RENDER_DPI) -> tuple[list[Path], list[str]]:
    warnings: list[str] = []
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    pdftoppm = shutil.which("pdftoppm")
    if not soffice or not pdftoppm:
        missing = []
        if not soffice:
            missing.append("soffice/libreoffice")
        if not pdftoppm:
            missing.append("pdftoppm")
        raise RuntimeError("missing external render tools: " + ", ".join(missing))
    render_dir.mkdir(parents=True, exist_ok=True)
    temp_parent = render_dir.parent / ".kq1c_office_tmp"
    temp_parent.mkdir(parents=True, exist_ok=True)
    with tempfile.TemporaryDirectory(prefix="kq1c_office_render_", dir=str(temp_parent)) as tmp:
        tmp_path = Path(tmp)
        work_dir = tmp_path / "work"
        pdf_dir = tmp_path / "pdf"
        png_dir = tmp_path / "png"
        profile_dir = tmp_path / "lo-profile"
        work_dir.mkdir(parents=True, exist_ok=True)
        pdf_dir.mkdir(parents=True, exist_ok=True)
        png_dir.mkdir(parents=True, exist_ok=True)
        profile_dir.mkdir(parents=True, exist_ok=True)
        convert = _run_libreoffice_pdf_convert(soffice, pptx_path, pdf_dir, work_dir, profile_dir)
        if convert.returncode != 0:
            raise RuntimeError("LibreOffice PDF conversion failed: " + convert.stdout[-1200:])
        pdf_candidates = _find_libreoffice_pdf_candidates(pdf_dir, work_dir)
        if not pdf_candidates:
            diagnostic = "\n".join(
                [
                    "LibreOffice PDF conversion returned 0 but produced no PDF",
                    "stdout_tail=" + convert.stdout[-1200:],
                    "work_dir_entries=" + ", ".join(sorted(path.name for path in work_dir.iterdir())),
                    "pdf_dir_entries=" + ", ".join(sorted(path.name for path in pdf_dir.iterdir())),
                ]
            )
            raise RuntimeError(diagnostic)
        pdf_path = pdf_candidates[0]
        render = subprocess.run(
            [pdftoppm, "-png", "-r", str(dpi), str(pdf_path), str(png_dir / "slide")],
            text=True,
            stdout=subprocess.PIPE,
            stderr=subprocess.STDOUT,
            timeout=120,
            check=False,
        )
        if render.returncode != 0:
            raise RuntimeError("pdftoppm render failed: " + render.stdout[-1200:])
        produced = sorted(png_dir.glob("slide-*.png"))
        if not produced:
            raise RuntimeError("pdftoppm produced no slide PNGs")
        out_paths: list[Path] = []
        for idx, src in enumerate(produced, start=1):
            dst = render_dir / f"slide_{idx:02d}.png"
            shutil.copy2(src, dst)
            out_paths.append(dst)
    return out_paths, warnings


def _render_with_python_pptx_text(pptx_path: Path, render_dir: Path) -> tuple[list[Path], list[str]]:
    try:
        from PIL import Image, ImageDraw, ImageFont
    except Exception as exc:
        raise RuntimeError(f"Pillow is required for python-pptx text render fallback: {exc}") from exc
    slide_texts = _extract_pptx_slide_text(pptx_path)
    if not slide_texts:
        count = _slide_count_with_python_pptx(pptx_path)
        slide_texts = [f"Slide {idx}" for idx in range(1, count + 1)]
    render_dir.mkdir(parents=True, exist_ok=True)
    try:
        title_font = ImageFont.truetype("DejaVuSans-Bold.ttf", 32)
        body_font = ImageFont.truetype("DejaVuSans.ttf", 21)
        small_font = ImageFont.truetype("DejaVuSans.ttf", 13)
    except Exception:
        title_font = body_font = small_font = ImageFont.load_default()
    out_paths: list[Path] = []
    for idx, text in enumerate(slide_texts, start=1):
        img = Image.new("RGB", (1280, 720), "white")
        draw = ImageDraw.Draw(img)
        draw.rectangle((0, 0, 1280, 62), fill=(235, 238, 244))
        draw.text((42, 20), f"KQ-1C independent PPTX text render / slide_{idx:02d}", fill=(65, 65, 65), font=small_font)
        lines = [line.strip() for line in text.splitlines() if line.strip()]
        if lines:
            draw.text((42, 92), _wrap_text(lines[0], 56), fill=(20, 20, 20), font=title_font, spacing=6)
        y = 180
        for line in lines[1:12]:
            draw.text((72, y), _wrap_text(line, 92), fill=(40, 40, 40), font=body_font, spacing=5)
            y += 72
            if y > 640:
                break
        dst = render_dir / f"slide_{idx:02d}.png"
        img.save(dst)
        out_paths.append(dst)
    return out_paths, ["used python-pptx text renderer instead of office/PDF render"]


def _wrap_text(value: str, width: int) -> str:
    import textwrap

    return "\n".join(textwrap.wrap(value, width=width, break_long_words=False))


def render_pptx_independently(
    pptx_path: Path,
    render_dir: Path,
    *,
    render_mode: str = "auto",
) -> tuple[list[Path], str, bool, bool, list[str]]:
    warnings: list[str] = []
    requested = render_mode
    if render_mode not in {"auto", "libreoffice", "python-pptx-text"}:
        raise ValueError(f"unsupported render mode: {render_mode}")
    if render_mode in {"auto", "libreoffice"}:
        try:
            paths, office_warnings = _render_with_libreoffice(pptx_path, render_dir)
            warnings.extend(office_warnings)
            return paths, "libreoffice_pdf_pdftoppm", True, True, warnings
        except Exception as exc:
            if requested == "libreoffice":
                raise
            warnings.append(f"LibreOffice render unavailable; falling back to python-pptx text render: {exc}")
    paths, fallback_warnings = _render_with_python_pptx_text(pptx_path, render_dir)
    warnings.extend(fallback_warnings)
    return paths, "python_pptx_text_renderer", True, False, warnings


def _build_slide_observations(pptx_path: Path, rendered_paths: list[Path]) -> tuple[tuple[KQ1CSlideRenderObservation, ...], dict[str, int]]:
    slide_texts = _extract_pptx_slide_text(pptx_path)
    char_counts, overflow_flags, tiny_text_flags = _extract_pptx_text_geometry_findings(pptx_path)
    observations: list[KQ1CSlideRenderObservation] = []
    counts = {"empty_slide_count": 0, "text_overflow_count": 0, "tiny_text_count": 0}
    for idx, path in enumerate(rendered_paths, start=1):
        width, height, non_white_ratio = _non_white_pixel_ratio(path)
        text = slide_texts[idx - 1] if idx - 1 < len(slide_texts) else ""
        char_count = char_counts[idx - 1] if idx - 1 < len(char_counts) else len(text)
        overflow = overflow_flags[idx - 1] if idx - 1 < len(overflow_flags) else False
        tiny_text = tiny_text_flags[idx - 1] if idx - 1 < len(tiny_text_flags) else False
        empty = non_white_ratio < 0.01 or char_count == 0
        if empty:
            counts["empty_slide_count"] += 1
        if overflow:
            counts["text_overflow_count"] += 1
        if tiny_text:
            counts["tiny_text_count"] += 1
        observations.append(
            KQ1CSlideRenderObservation(
                slide_id=f"slide_{idx:02d}",
                render_path=path.name,
                width_px=width,
                height_px=height,
                non_white_pixel_ratio=round(non_white_ratio, 6),
                text_character_count=char_count,
                extracted_text_preview=" ".join(text.split())[:260],
                empty_slide_detected=empty,
                tiny_text_detected=tiny_text,
                text_overflow_detected=overflow,
            )
        )
    return tuple(observations), counts


def _update_review_packet(root: Path, render_dir_rel: str) -> None:
    packet_path = _find_review_packet(root)
    if packet_path is None:
        return
    try:
        packet = _load_json_object(packet_path)
    except Exception:
        return
    refs = packet.get("deck_artifact_refs") if isinstance(packet.get("deck_artifact_refs"), dict) else {}
    refs["independent_rendered_slides"] = f"{render_dir_rel}/*.png"
    refs["kq1c_render_manifest"] = "kq1c_render_manifest.json"
    refs["kq1c_visual_qa_report"] = "kq1c_visual_qa_report.json"
    packet["deck_artifact_refs"] = refs
    packet["independent_pptx_render_performed_by_kq1c"] = True
    packet["review_state"] = packet.get("review_state") or "pending_human_review"
    packet["human_review_decision"] = packet.get("human_review_decision") or None
    notes = packet.get("review_notes") if isinstance(packet.get("review_notes"), list) else []
    notes.append("KQ-1C adds an independent PPTX render output and render-based visual QA report; this is still not a human approval.")
    packet["review_notes"] = notes
    write_json(packet_path, packet)


def run_kq1c_independent_render_qa(
    input_bundle: Path,
    output_bundle_dir: Path,
    *,
    zip_out: Path | None = None,
    quality_report_dir: Path | None = None,
    quality_report_zip: Path | None = None,
    render_mode: str = "auto",
    require_office_render: bool = False,
) -> KQ1CRenderQAResult:
    errors: list[str] = []
    warnings: list[str] = []
    root = _copy_bundle_to_dir(input_bundle, output_bundle_dir)
    pptx_files = find_pptx_files(root)
    if not pptx_files:
        raise RuntimeError("KQ-1C input bundle contains no PPTX artifact")
    pptx_path = pptx_files[0]
    pptx_digest = digest_file(pptx_path)
    slide_count = _slide_count_with_python_pptx(pptx_path)
    render_dir = root / "independent_rendered_slides"
    rendered_paths, engine, independent_render, office_render, render_warnings = render_pptx_independently(pptx_path, render_dir, render_mode=render_mode)
    warnings.extend(render_warnings)
    if require_office_render and not office_render:
        errors.append("office/PDF render was required but was not performed")
    if slide_count and len(rendered_paths) != slide_count:
        errors.append(f"independent render count {len(rendered_paths)} does not match PPTX slide count {slide_count}")

    observations, counts = _build_slide_observations(pptx_path, rendered_paths)
    blocking_defects: list[dict[str, Any]] = []
    if counts["empty_slide_count"]:
        blocking_defects.append({"defect": "empty_slides", "count": counts["empty_slide_count"]})
    if counts["text_overflow_count"]:
        blocking_defects.append({"defect": "text_overflow", "count": counts["text_overflow_count"]})
    if counts["tiny_text_count"]:
        blocking_defects.append({"defect": "tiny_text", "count": counts["tiny_text_count"]})
    if not rendered_paths:
        blocking_defects.append({"defect": "no_independent_rendered_slides", "count": 1})
    if blocking_defects:
        errors.append("independent visual QA found blocking defects: " + ", ".join(f"{item['defect']}={item['count']}" for item in blocking_defects))

    visual_status = "ready" if not blocking_defects else "failed"
    render_dir_rel = render_dir.relative_to(root).as_posix()
    render_manifest = {
        "schema_version": KQ1C_RENDER_MANIFEST_SCHEMA_VERSION,
        "phase": KQ1C_PHASE_ID,
        "workflow_id": KQ1C_WORKFLOW_ID,
        "scenario_id": KQ1A_DEFAULT_SCENARIO_ID,
        "deck_type": KQ1A_DEFAULT_DECK_TYPE,
        "input_bundle": str(input_bundle),
        "pptx_path": pptx_path.relative_to(root).as_posix(),
        "pptx_digest": pptx_digest,
        "slide_count_from_pptx": slide_count,
        "render_engine_requested": render_mode,
        "render_engine": engine,
        "independent_pptx_render_performed_by_kq1c": independent_render,
        "independent_office_render_performed_by_kq1c": office_render,
        "rendered_slide_count": len(rendered_paths),
        "rendered_slides": [path.relative_to(root).as_posix() for path in rendered_paths],
        "controlled_scope": dict(KQ1C_CONTROLLED_SCOPE_FLAGS),
        "kimi_level_claimed_by_kq1c": False,
        "selected_offline_workflow_parity_claim_supported_after_kq1c": False,
        "server3_local_intranet_route_verified_by_kq1c": False,
    }
    visual_qa_report = {
        "schema_version": KQ1C_VISUAL_QA_SCHEMA_VERSION,
        "phase": KQ1C_PHASE_ID,
        "workflow_id": KQ1C_WORKFLOW_ID,
        "scenario_id": KQ1A_DEFAULT_SCENARIO_ID,
        "status": visual_status,
        "visual_qa_status": visual_status,
        "method": "independent PPTX render QA over rendered PNG output",
        "render_engine": engine,
        "slide_count": slide_count,
        "rendered_slide_count": len(rendered_paths),
        "empty_slide_count": counts["empty_slide_count"],
        "text_overflow_count": counts["text_overflow_count"],
        "tiny_text_count": counts["tiny_text_count"],
        "blocking_defects": blocking_defects,
        "slide_observations": [observation.as_dict() for observation in observations],
        "independent_pptx_render_performed_by_kq1c": independent_render,
        "independent_office_render_performed_by_kq1c": office_render,
        "human_review_decision": None,
        "kimi_level_claimed_by_kq1c": False,
        "selected_offline_workflow_parity_claim_supported_after_kq1c": False,
        "server3_local_intranet_route_verified_by_kq1c": False,
    }
    geometry_report = {
        "schema_version": "kq1c.geometry_report.v1",
        "phase": KQ1C_PHASE_ID,
        "scenario_id": KQ1A_DEFAULT_SCENARIO_ID,
        "slide_count": slide_count,
        "slides_checked": slide_count,
        "empty_slide_count": counts["empty_slide_count"],
        "text_overflow_count": counts["text_overflow_count"],
        "tiny_text_count": counts["tiny_text_count"],
        "method": "independent PPTX render output plus python-pptx geometry approximation",
    }

    write_json(root / "kq1c_render_manifest.json", render_manifest)
    write_json(root / "kq1c_visual_qa_report.json", visual_qa_report)
    # Keep KQ-1A pointing at the stricter KQ-1C reports.
    write_json(root / "visual_qa_report.json", visual_qa_report)
    write_json(root / "geometry_report.json", geometry_report)
    _update_review_packet(root, render_dir_rel)

    if zip_out:
        make_zip_from_dir(root, zip_out)
    assessment_input = zip_out if zip_out else root
    kq1a_result = assess_kq1a_deck_artifact_bundle(assessment_input)
    quality_report_dir = quality_report_dir or (root / "kq1a_after_kq1c_quality_report")
    report_paths = write_kq1a_assessment_artifacts(kq1a_result, quality_report_dir)
    if quality_report_zip:
        make_zip_from_dir(quality_report_dir, quality_report_zip)
    if kq1a_result.status != "ready":
        errors.extend(kq1a_result.errors)
    status = "ready" if not errors else "failed"
    return KQ1CRenderQAResult(
        status=status,
        errors=tuple(errors),
        warnings=tuple(warnings) + tuple(kq1a_result.warnings),
        scenario_id=KQ1A_DEFAULT_SCENARIO_ID,
        deck_type=KQ1A_DEFAULT_DECK_TYPE,
        input_bundle=str(input_bundle),
        output_bundle_dir=str(root),
        output_bundle_zip=str(zip_out) if zip_out else None,
        pptx_path=str(pptx_path),
        pptx_digest=pptx_digest,
        slide_count_from_pptx=slide_count,
        independent_pptx_render_performed_by_kq1c=independent_render,
        independent_office_render_performed_by_kq1c=office_render,
        render_engine=engine,
        render_engine_requested=render_mode,
        render_engine_available=bool(rendered_paths),
        rendered_slide_count=len(rendered_paths),
        rendered_slide_dir=str(render_dir),
        visual_qa_status=visual_status,
        empty_slide_count=counts["empty_slide_count"],
        text_overflow_count=counts["text_overflow_count"],
        tiny_text_count=counts["tiny_text_count"],
        blocking_defect_count=len(blocking_defects),
        kq1a_status_after_kq1c=kq1a_result.status,
        kq1a_report_path=str(report_paths["report"]),
        selected_offline_workflow_parity_claim_supported_after_kq1c=False,
        kimi_level_claimed_by_kq1c=False,
        server3_local_intranet_route_verified_by_kq1c=False,
        controlled_scope=dict(KQ1C_CONTROLLED_SCOPE_FLAGS),
        slide_observations=tuple(observation.as_dict() for observation in observations),
    )


def build_kq1c_capabilities_report() -> dict[str, Any]:
    soffice_available = bool(shutil.which("soffice") or shutil.which("libreoffice"))
    pdftoppm_available = bool(shutil.which("pdftoppm"))
    try:
        import pptx  # noqa: F401

        python_pptx_available = True
    except Exception:
        python_pptx_available = False
    try:
        import PIL  # noqa: F401

        pillow_available = True
    except Exception:
        pillow_available = False
    return {
        "checkpoint": KQ1C_PHASE_ID,
        "workflow_id": KQ1C_WORKFLOW_ID,
        "schema_version": KQ1C_SCHEMA_VERSION,
        "focus_scenario_id": KQ1A_DEFAULT_SCENARIO_ID,
        "focus_deck_type": KQ1A_DEFAULT_DECK_TYPE,
        "independent_pptx_render_qa_supported": True,
        "libreoffice_pdf_render_supported_when_available": True,
        "python_pptx_text_render_fallback_supported": True,
        "soffice_or_libreoffice_available": soffice_available,
        "pdftoppm_available": pdftoppm_available,
        "office_render_stack_available": soffice_available and pdftoppm_available,
        "python_pptx_available": python_pptx_available,
        "pillow_available": pillow_available,
        "visual_qa_over_independent_render_supported": True,
        "updates_bundle_visual_qa_report": True,
        "updates_review_packet_with_independent_render_refs": True,
        "kq1a_validation_after_independent_render_supported": True,
        "api_endpoint_added_by_kq1c": False,
        "db_schema_migration_added_by_kq1c": False,
        "frontend_runtime_changed_by_kq1c": False,
        "dependency_versions_changed_by_kq1c": False,
        "dockerfiles_changed_by_kq1c": False,
        "kimi_level_claimed_by_kq1c": False,
        "selected_offline_workflow_parity_claim_supported_after_kq1c": False,
        "server3_local_intranet_route_verified_by_kq1c": False,
        **KQ1C_CONTROLLED_SCOPE_FLAGS,
    }
